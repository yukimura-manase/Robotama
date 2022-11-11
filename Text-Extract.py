
### < Word, Excel, PowerPoint, PDF, TXT, CSV から文字列を抽出する方法: 関数のSample🔥 >

# 1. ElasticSearch (全文検索エンジン)に文章データを解析処理した後に、登録するために、Pythonを使用した時に、実際に使った関数たちです！

# => 解析の対象ファイルたちは、Word, Excel, PowerPoint, PDF, TXT, CSV

# 5つの関数で対応 => TXT, CSV は、同じ関数で対応可能でした。


# -------------------------------------------- Text-Extract-Module の領域展開 --------------------------------------------

# Word
import docx

# Excel
import openpyxl

# PowerPoint
from pptx import Presentation

# PDF
from pdfminer.layout import LAParams, LTContainer, LTTextBox
from pdfminer.pdfinterp import PDFPageInterpreter, PDFResourceManager
from pdfminer.converter import TextConverter
from pdfminer.pdfpage import PDFPage


from io import StringIO


# chardetによるエンコーディングの判定とテキストデータのデコード
# sjis.txtの内容：このファイルはシフトJISでエンコーディングされています
from chardet import detect


# -------------------------------------------- Text-Extract-Function の領域展開 --------------------------------------------


# TXT or CSVファイルのテキストデータ抽出・関数
def txt_or_csv_parse_func(file):

    with open(file, 'rb') as f:  # バイナリファイルとしてファイルをオープン

        b = f.read()  # ファイルの内容を全て読み込む

    enc = detect(b)  # chardet.detect関数を使ってエンコーディングを判定

    # 得られたエンコーディング情報を使ってファイルをオープンし直す
    with open(file, encoding=enc['encoding']) as f:
        s = f.read()


    # もしくは得られたエンコーディング情報を使ってバイト列をデコード
    s = b.decode(encoding=enc['encoding'])

    return repr(s)



# Wordテキストデータ抽出・関数
def get_all_text_from_docx(filepath: str):
    """
    Wordファイルから全文字列を取得する
    """
    document = docx.Document(filepath)
    return "\n".join(list(map(lambda par: par.text, document.paragraphs)))



# Excelテキストデータ抽出・関数
def excel_parse_func(file):

    # ブックを取得
    book = openpyxl.load_workbook(file)

    # print(book)

    # シート名を配列で取得
    sheet_name_list = book.sheetnames

    # print(sheet_name_list)

    results = []

    for sheet_name in sheet_name_list:
        # シート名を1つずつ取得 => sheetインスタンス
        sheet = book[sheet_name]
        # print(sheet)

        cell_list = [f"{cell.value}" for cells in tuple(sheet.columns) for cell in cells]
        results.extend(cell_list)

    return results



# PowerPointテキストデータ抽出・関数
def get_all_text_from_pptx(filepath: str):
    """
    PoerPointファイルの文字列を取得する
    """
    presentation = Presentation(filepath)
    results = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            # 文字列以外は飛ばす
            if not shape.has_text_frame:
                continue
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    results.append(run.text)
    return results


# PDFテキストデータ抽出・関数
def get_all_text_from_pdf(filepath: str):
    rsrcmgr = PDFResourceManager()
    outfp = StringIO()
    laparams = LAParams()
    laparams.detect_vertical = True
    device = TextConverter(rsrcmgr, outfp, codec='utf-8', laparams=laparams)
    interpreter = PDFPageInterpreter(rsrcmgr, device)

    with open(filepath, "rb") as file:
        for page in PDFPage.get_pages(file, pagenos=None, maxpages=0, password=None, caching=True, check_extractable=True):
            interpreter.process_page(page)

    ret = outfp.getvalue()
    # 後始末をしておく
    device.close()
    outfp.close()
    return ret.splitlines()




# ------------------------------------------------------------------------------------------------------------------------------------------------















